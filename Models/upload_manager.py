import json
import multiprocessing
import os
import re
import subprocess
import traceback
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path
from typing import List, Optional, Dict, Any
import eyed3
import fitz
import requests
import spotipy
import unicodedata
import yaml
from PIL import Image, ExifTags
from PyQt6.QtCore import QObject, QThread, pyqtSignal, QTimer
from docx import Document
from mutagen import File
from mutagen.flac import Picture
from mutagen.id3 import APIC
from mutagen.mp4 import MP4Cover
from pptx import Presentation
from pyOneFichierClient.OneFichierAPI.exceptions import FichierResponseNotOk
from pyrebase import pyrebase
from requests_toolbelt.multipart.encoder import MultipartEncoder, MultipartEncoderMonitor
from spotipy.oauth2 import SpotifyClientCredentials

from Core.settings import FFPROBE_PATH, FFMPEG_PATH, BASE_DIR
from Models.category import CatManager
from Service.py1FichierClient import FichierClient, s

AFF_ID = "5091183"


def sanitize_for_firebase_key(key: str) -> str:
    if not key:
        return "untitled"
    key = unicodedata.normalize('NFKD', key).encode('ascii', 'ignore').decode('ascii')
    key = re.sub(r'[.#$/\[\]()]+', '_', key)
    key = key.replace(' ', '_')
    key = re.sub(r'_+', '_', key)
    return key.strip('_')


def sanitize_folder_name(name: str) -> str:
    name = unicodedata.normalize('NFKD', name).encode('ascii', 'ignore').decode('ascii')
    name = re.sub(r'[^\w\-]', '_', name)
    name = re.sub(r'_+', '_', name)
    return name.strip('_')


def sanitize_dict_keys(d):
    if isinstance(d, dict):
        new_dict = {}
        for k, v in d.items():
            new_key = sanitize_for_firebase_key(str(k))
            new_dict[new_key] = sanitize_dict_keys(v)
        return new_dict
    elif isinstance(d, list):
        return [sanitize_dict_keys(i) for i in d]
    else:
        return d


def extract_audio_cover(file_path: str, output_jpg: str) -> bool:
    try:
        audio = File(file_path)
        if audio is None:
            return False

        if audio.tags is not None:
            for tag in audio.tags.values():
                if isinstance(tag, APIC):
                    with open(output_jpg, 'wb') as img:
                        img.write(tag.data)
                    return True

        if hasattr(audio, "pictures"):
            for pic in audio.pictures:
                if isinstance(pic, Picture):
                    with open(output_jpg, 'wb') as img:
                        img.write(pic.data)
                    return True

        if hasattr(audio, "tags") and 'covr' in audio.tags:
            for cover in audio.tags['covr']:
                if isinstance(cover, MP4Cover):
                    with open(output_jpg, 'wb') as img:
                        img.write(cover)
                    return True
    except Exception as e:
        print(f"Erreur extraction cover audio : {e}")
    return False


def generate_thumbnail(video_path: str, thumb_path: str, percent: float = 0.15):
    cmd_probe = [str(FFPROBE_PATH), '-v', 'error', '-show_entries', 'format=duration', '-of',
                 'default=noprint_wrappers=1:nokey=1', video_path]
    proc = subprocess.run(cmd_probe, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffprobe error: {proc.stderr}")
    duration_str = proc.stdout.strip()
    if not duration_str or not duration_str.replace('.', '', 1).isdigit():
        raise ValueError(f"Durée invalide extraite: '{duration_str}'")
    duration = float(duration_str)

    timestamp = duration * percent
    filt = "crop='if(gt(iw/ih\\,16/9)\\,ih*16/9\\,iw)':if(gt(iw/ih\\,16/9)\\,ih\\,iw*9/16),scale=320:180"
    cmd_ff = [str(FFMPEG_PATH), '-hide_banner', '-loglevel', 'error', '-hwaccel', 'auto', '-ss', str(timestamp), '-i',
              video_path, '-vf', filt, '-vframes', '1', '-qscale:v', '2', '-preset', 'ultrafast', '-threads',
              str(multiprocessing.cpu_count()), '-nostdin', '-y', thumb_path]
    proc2 = subprocess.run(cmd_ff, capture_output=True, text=True)
    if proc2.returncode != 0:
        raise RuntimeError(f"ffmpeg error: {proc2.stderr}")
    return thumb_path


def safe_serialize(obj):
    """Convertit les objets non sérialisables en strings"""
    if isinstance(obj, (str, int, float, bool, type(None))):
        return obj
    elif isinstance(obj, (list, tuple)):
        return [safe_serialize(item) for item in obj]
    elif isinstance(obj, dict):
        return {str(k): safe_serialize(v) for k, v in obj.items()}
    elif hasattr(obj, '__str__'):
        return str(obj)
    else:
        return repr(obj)


class TMDBClient:
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.base_url = "https://api.themoviedb.org/3"

    def search_movie(self, query: str, year: Optional[int] = None) -> Optional[Dict]:
        url = f"{self.base_url}/search/movie"
        params = {"api_key": self.api_key, "query": query}
        if year:
            params["year"] = year
        response = requests.get(url, params=params)
        if response.status_code == 200:
            return response.json().get("results", [])[0] if response.json().get("results") else None
        return None

    def search_tv(self, query: str, year: Optional[int] = None) -> Optional[Dict]:
        url = f"{self.base_url}/search/tv"
        params = {"api_key": self.api_key, "query": query}
        if year:
            params["first_air_date_year"] = year
        response = requests.get(url, params=params)
        if response.status_code == 200:
            return response.json().get("results", [])[0] if response.json().get("results") else None
        return None

    def get_tv_details(self, tv_id: int) -> Optional[Dict]:
        url = f"{self.base_url}/tv/{tv_id}"
        params = {"api_key": self.api_key}
        response = requests.get(url, params=params)
        return response.json() if response.status_code == 200 else None

    def get_movie_details(self, movie_id: int) -> Optional[Dict]:
        url = f"{self.base_url}/movie/{movie_id}"
        params = {"api_key": self.api_key}
        response = requests.get(url, params=params)
        return response.json() if response.status_code == 200 else None


def parse_video_filename(filename: str) -> Dict[str, Any]:
    name = Path(filename).stem
    patterns = [
        r"^(.*?)[ ._-]s(\d{1,2})[ ._-]?e(\d{1,2})",
        r"^(.*?)[ ._-]\((\d{4})\)",
        r"^(.*?)[ ._-](\d{4})"
    ]

    for pattern in patterns:
        match = re.search(pattern, name, re.IGNORECASE)
        if match:
            if "s" in pattern and "e" in pattern:
                return {
                    "type": "tv",
                    "title": match.group(1).replace('.', ' ').strip(),
                    "season": int(match.group(2)),
                    "episode": int(match.group(3))
                }
            else:
                return {
                    "type": "movie",
                    "title": match.group(1).replace('.', ' ').strip(),
                    "year": int(match.group(2))
                }

    return {"type": "unknown", "title": name}


def extract_music_metadata(file_path: str) -> Dict[str, Any]:
    """Extrait les métadonnées musicales avec eyed3 comme dans le premier code"""
    metadata = {
        'artist': 'Artiste Inconnu',
        'title': '',
        'year': None
    }

    try:
        filename = os.path.splitext(os.path.basename(file_path))[0]
        audio = eyed3.load(file_path)
        if ' - ' in filename:
            artist, title_part = filename.split(' - ', 1)
            metadata['artist'] = artist.strip()
            metadata['title'] = re.sub(r'\(.*?\)', '', title_part).strip()
        else:
            metadata['title'] = filename

        if audio and audio.tag:
            tag = audio.tag
            if tag.artist:
                metadata['artist'] = tag.artist.strip()
            if tag.title:
                metadata['title'] = tag.title.strip()
            if tag.getBestDate():
                metadata['year'] = str(tag.getBestDate().year)

        metadata['title'] = re.sub(
            rf"^{re.escape(metadata['artist'])}\s*-\s*",
            "",
            metadata['title'],
            flags=re.IGNORECASE
        )

        metadata['title'] = re.sub(r'\(.*?\)|\[.*?\]', '', metadata['title']).strip()
        metadata['title'] = re.sub(
            r'\b(official music video|official video|lyrics|audio|remastered|remaster|version)\b',
            '', metadata['title'], flags=re.IGNORECASE
        ).strip()

    except Exception as e:
        print(f"Erreur extraction métadonnées musicales: {e}")
        filename = os.path.splitext(os.path.basename(file_path))[0]
        if ' - ' in filename:
            artist, title = filename.split(' - ', 1)
            metadata['artist'] = artist
            metadata['title'] = re.sub(r'\(.*?\)', '', title).strip()
        else:
            metadata['title'] = filename

    return metadata


class UploadThread(QThread):
    progress = pyqtSignal(int)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, client: FichierClient, file_path: str):
        super().__init__()
        self.client = client
        self.file_path = file_path
        self.last_logged_percent = -10

    def run(self):
        try:
            resp = self.client.api_call('https://api.1fichier.com/v1/upload/get_upload_server.cgi', method='GET')
            up_srv, upload_id = resp['url'], resp['id']
            with open(self.file_path, 'rb') as f:
                encoder = MultipartEncoder({'file[]': (Path(self.file_path).name, f, 'application/octet-stream')})

                def monitor_callback(mon):
                    percent = int(100 * mon.bytes_read / mon.len)
                    if percent >= self.last_logged_percent + 10:
                        self.last_logged_percent = (percent // 10) * 10
                        print(f"Progress: {self.last_logged_percent}%")
                    self.progress.emit(percent)

                monitor = MultipartEncoderMonitor(encoder, monitor_callback)

                headers = {'Content-Type': monitor.content_type}
                if self.client.authed:
                    headers.update(self.client.auth_nc)
                url = f'https://{up_srv}/upload.cgi?id={upload_id}'
                r = s.post(url, data=monitor, headers=headers, allow_redirects=False)
                if 'Location' not in r.headers:
                    raise FichierResponseNotOk('Header Location manquant')
                loc = r.headers['Location']
                r2 = s.get(f'https://{up_srv}{loc}')
                m = re.search(r'<td class="normal"><a href="(.+)"', r2.text)
                if not m:
                    raise FichierResponseNotOk('Lien de téléchargement introuvable')
                self.finished.emit(m.group(1))
        except Exception as e:
            tb = traceback.format_exc()
            self.error.emit(f"{e}\n{tb}")


class UploadManager(QObject):
    progress = pyqtSignal(int)
    file_progress = pyqtSignal(int)
    finished = pyqtSignal(str, str)
    thumb_finished = pyqtSignal(str, str)
    error = pyqtSignal(str)
    all_done = pyqtSignal()

    def __init__(
            self,
            client: FichierClient,
            cat_manager: CatManager,
            existing_files: dict,
            folder_ids_by_name: dict,
            root_folder_id: str = '0',
            local_path: str = "",
    ):
        super().__init__()
        self.thumb_threads: List[UploadThread] = []
        self._firebase_keys: dict[str, tuple[str, str, str]] = {}
        self.client = client
        self.local_path = local_path
        self.cat_manager = cat_manager
        self.existing_files = existing_files
        self.folder_ids_by_name = folder_ids_by_name
        self.folder_parent_id = root_folder_id
        with open(BASE_DIR / 'Config/config.yaml', 'r') as f:
            cfg = yaml.safe_load(f)
        firebase = pyrebase.initialize_app(cfg['firebase'])
        user = firebase.auth().sign_in_with_email_and_password(
            cfg.get('email', 'weck20pro@gmail.com'),
            cfg.get('password', 'Azerty')
        )
        self.db = firebase.database()
        self.token = user['idToken']
        self.uid = user['localId']
        self.files_to_upload: List[tuple[str, bool]] = []
        self.files_done = 0
        self.total_files = 0
        self.thread: UploadThread | None = None
        self.thumb_finished.connect(self._on_thumb_finished)
        self.refresh_token = user['refreshToken']
        self.token = user['idToken']
        self.auth = firebase.auth()
        self.current_file_progress = 0
        self.tmdb_metadata = {}
        self.music_metadata = {}
        self.tmdb_client = TMDBClient(cfg['tmdb']['api_key'])

        spotify_config = cfg.get('spotify', {})
        if spotify_config.get('client_id') and spotify_config.get('client_secret'):
            auth_manager = SpotifyClientCredentials(
                client_id=spotify_config['client_id'],
                client_secret=spotify_config['client_secret']
            )
            self.spotify_client = spotipy.Spotify(auth_manager=auth_manager)
        else:
            self.spotify_client = None

    def _ensure_token(self):
        try:
            new_user = self.auth.refresh(self.refresh_token)
            self.token = new_user['idToken']
            self.refresh_token = new_user['refreshToken']
        except Exception as e:
            print("[Firebase] Impossible de rafraîchir le token :", e)

    def get_uploaded_files_from_firebase(self) -> dict[str, set[str]]:
        try:
            user_data = self.db.child('users').child(self.uid).get(self.token).val()
            uploaded = {}
            if user_data:
                for category_key, files in user_data.items():
                    if not isinstance(files, dict):
                        print(f"[Firebase] Warning: données inattendues pour catégorie '{category_key}': {files}")
                        continue
                    uploaded[category_key.lower()] = set(file_key.lower() for file_key in files.keys())
            return uploaded
        except Exception as e:
            print(f"[Firebase] Erreur lecture fichiers uploadés : {e}")
            return {}

    def _get_db_ref(self, category: str, title: str):
        self._ensure_token()
        safe_category = sanitize_for_firebase_key(category)
        safe_title = sanitize_for_firebase_key(title)
        return self.db.child('users').child(self.uid).child(safe_category).child(safe_title)

    def store_metadata_in_firebase(
            self, main_link: str,
            metadata: dict,
            thumb_link: str = '',
            category: str = None,
            title: str = None,
            file_extension: str = None,
            music_metadata: str = None,
            local_path: str = None  # Ajouter local_path
    ) -> bool:
        self._ensure_token()
        clean_metadata = sanitize_dict_keys(metadata)
        clean_metadata = safe_serialize(clean_metadata)

        data = {
            'file_link': main_link,
            'metadata': clean_metadata,
            'thumbnail_link': thumb_link,
            'file_extension': file_extension
        }

        if music_metadata:
            data['music_metadata'] = music_metadata

        if local_path:
            data['local_path'] = local_path  # Ajouter local_path

        try:
            self._get_db_ref(category, title).set(data, self.token)
            return True
        except Exception as e:
            print(f"[Firebase] Erreur stockage métadonnées : {e}")
            return False

    def set_files(self, files: List[tuple[str, bool]]):
        self.files_to_upload = files
        self.files_done = 0
        self.total_files = len(files)
        self.existing_files = self.get_uploaded_files_from_firebase()

    def start(self):
        if not self.files_to_upload:
            self.all_done.emit()
        else:
            self.upload_next_file()

    def _find_album_with_track(self, artist_id: str, track_name: str) -> Dict[str, Any]:
        """Parcourt tous les albums d'un artiste pour trouver celui qui contient un morceau spécifique"""
        if not self.spotify_client:
            return {}

        try:
            albums = []
            offset = 0
            while True:
                batch = self.spotify_client.artist_albums(
                    artist_id,
                    album_type='album,single,compilation',
                    limit=50,
                    offset=offset
                )
                albums.extend(batch['items'])
                if not batch['next']:
                    break
                offset += 50

            normalized_track_name = self.normalize_string(track_name)

            for album in albums:
                try:
                    tracks = self.spotify_client.album_tracks(album['id'])['items']
                    for track in tracks:
                        if self.normalize_string(track['name']) == normalized_track_name:
                            return album
                except Exception as e:
                    print(f"Erreur lors de la récupération des pistes de l'album {album['name']}: {e}")
                    continue

        except Exception as e:
            print(f"Erreur lors de la recherche d'album: {e}")

        return {}

    def upload_next_file(self):
        if not self.files_to_upload:
            self.progress.emit(100)
            self.all_done.emit()
            return

        self.current_file_progress = 0

        path, _ = self.files_to_upload.pop(0)
        suffix = Path(path).suffix.lower()
        category = self.cat_manager.get_category(suffix).capitalize()
        local_path = self.local_path
        key = category.lower()

        name = Path(path).stem
        metadata = self.get_all_metadata(path, category)

        title_from_meta = name
        if key == 'videos':
            tags = metadata.get('ffprobe', {}).get('format', {}).get('tags', {})
            if tags.get('title', '').strip():
                title_from_meta = tags['title'].strip()
            elif key == 'musiques':
                clean_title = title_from_meta
                clean_title = re.sub(
                    r'(official.*video|music video|lyric video|audio|official|videoclip|clip officiel)',
                    '', clean_title, flags=re.IGNORECASE)
                clean_title = re.sub(r'[\(\[].*?[\)\]]', '', clean_title)
                clean_title = re.sub(r'\s+', ' ', clean_title).strip()
                clean_title = re.sub(r'^[-\s]+|[-\s]+$', '', clean_title)

                spotify_meta = metadata.get('spotify', {})
                basic_meta = metadata

                if spotify_meta and spotify_meta.get('spotify_artist') and spotify_meta.get('spotify_album'):
                    artist = spotify_meta['spotify_artist'][0]
                    album = spotify_meta['spotify_album']
                    track_title = basic_meta.get('title', clean_title)
                    new_title = f"{artist} - {track_title}"

                elif basic_meta.get('artist') and basic_meta.get('title'):
                    new_title = f"{basic_meta['artist']} - {basic_meta['title']}"

                else:
                    new_title = clean_title

        firebase_sanitized_title = sanitize_folder_name(title_from_meta).lower()

        if firebase_sanitized_title in self.existing_files.get(key, set()):
            self.files_done += 1
            self.progress.emit(int(self.files_done / self.total_files * 100))
            QTimer.singleShot(0, self.upload_next_file)
            return

        if key not in self.folder_ids_by_name:
            safe_name = sanitize_folder_name(key)

            resp = self.client.create_folder(safe_name, parent_folder_id=self.folder_parent_id)
            self.folder_ids_by_name[key] = resp.get('folder_id') or resp.get('id')

        is_media = key in ('musiques', 'videos')

        self.thread = UploadThread(self.client, path)
        self.thread.progress.connect(self.on_file_progress)
        self.thread.finished.connect(lambda link, p=path: self.on_finished(link, p, category, name, is_media, local_path))
        self.thread.error.connect(self.error.emit)
        self.thread.start()

    def on_file_progress(self, percent: int):
        self.current_file_progress = percent
        if self.total_files > 0:
            global_progress = int(
                (self.files_done / self.total_files * 100) +
                (percent / (100 * self.total_files))
            )
            self.progress.emit(global_progress)
            self.file_progress.emit(percent)

    @staticmethod
    def render_first_page_a4(doc_path: str, dpi: int = 150) -> tuple[bytes, str] | tuple[None, None]:
        mm_to_inch = 25.4
        a4_w_in = 210 / mm_to_inch
        a4_h_in = 297 / mm_to_inch
        a4_w_px = int(a4_w_in * dpi)
        a4_h_px = int(a4_h_in * dpi)

        ext = Path(doc_path).suffix.lower()
        pil_img = None

        if ext == '.pdf':
            doc = fitz.open(str(doc_path))
            if doc.page_count:
                page = doc.load_page(0)
                m = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=m, alpha=False)
                pil_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            doc.close()

        elif ext == '.pptx':
            from pptx import Presentation
            pres = Presentation(str(doc_path))
            if pres.slides:
                slide = pres.slides[0]
                tmp_pdf = BytesIO()
                pres.save(tmp_pdf)
                tmp_pdf.seek(0)
                doc = fitz.open(stream=tmp_pdf.read(), filetype="pdf")
                page = doc.load_page(0)
                m = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=m, alpha=False)
                pil_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                doc.close()

        elif ext == '.docx':
            from docx import Document
            docx = Document(str(doc_path))
            for rel in docx.part._rels.values():
                if 'image' in rel.reltype:
                    data = rel.target_part.blob
                    pil_img = Image.open(BytesIO(data))
                    break

        if pil_img is None:
            return None, None

        img = pil_img.convert('RGB')
        img.thumbnail((a4_w_px, a4_h_px), Image.Resampling.LANCZOS)
        thumb = Image.new('RGB', (a4_w_px, a4_h_px), (255, 255, 255))
        x = (a4_w_px - img.width) // 2
        y = (a4_h_px - img.height) // 2
        thumb.paste(img, (x, y))

        buf = BytesIO()
        thumb.save(buf, 'JPEG', quality=90)
        return buf.getvalue(), '.jpg'

    def get_tmdb_metadata(self, file_path: str) -> Dict[str, Any]:
        filename = os.path.basename(file_path)
        parsed_info = parse_video_filename(filename)

        if parsed_info["type"] == "movie":
            movie = self.tmdb_client.search_movie(parsed_info["title"], parsed_info.get("year"))
            if movie:
                details = self.tmdb_client.get_movie_details(movie["id"])
                return {
                    "type": "movie",
                    "title": details.get("title"),
                    "original_title": details.get("original_title"),
                    "overview": details.get("overview"),
                    "genres": [genre["name"] for genre in details.get("genres", [])],
                    "release_date": details.get("release_date"),
                    "tmdb_id": details.get("id"),
                    "poster_path": details.get("poster_path"),
                    "backdrop_path": details.get("backdrop_path")
                }

        elif parsed_info["type"] == "tv":
            tv_show = self.tmdb_client.search_tv(parsed_info["title"], parsed_info.get("year"))
            if tv_show:
                details = self.tmdb_client.get_tv_details(tv_show["id"])
                return {
                    "type": "tv",
                    "title": details.get("name"),
                    "original_title": details.get("original_name"),
                    "overview": details.get("overview"),
                    "genres": [genre["name"] for genre in details.get("genres", [])],
                    "first_air_date": details.get("first_air_date"),
                    "season": parsed_info.get("season"),
                    "episode": parsed_info.get("episode"),
                    "tmdb_id": details.get("id"),
                    "poster_path": details.get("poster_path"),
                    "backdrop_path": details.get("backdrop_path")
                }

        return {"type": "unknown"}

    def get_spotify_metadata(self, artist: str, track: str = None, album: str = None) -> Dict[str, Any]:
        if not self.spotify_client:
            return {}

        try:
            filename = f"{artist} - {track}" if artist and track else artist
            metadata = {
                'artist': artist or 'Artiste Inconnu',
                'album': album or 'Album Inconnu',
                'title': track or '',
                'genre': 'Inconnu',
                'year': None
            }
            clean = re.sub(r'\(.*?\)|\[.*?\]', '', metadata['title']).strip()
            clean = re.sub(
                r'\b(official music video|official video|lyrics|audio|remastered|remaster|version)\b',
                '', clean, flags=re.IGNORECASE
            ).strip()

            if metadata['artist'] != 'Artiste Inconnu':
                q = f'artist:"{metadata["artist"]}" track:"{clean}"'
                res = self.spotify_client.search(q=q, type='track', limit=10)
                items = res['tracks']['items']

                best, score = None, 0
                for tr in items:
                    t_title = tr['name'].lower()
                    t_artist = tr['artists'][0]['name'].lower()

                    s = (
                            SequenceMatcher(None, clean.lower(), t_title).ratio() * 0.7 +
                            SequenceMatcher(None, metadata['artist'].lower(), t_artist).ratio() * 0.3
                    )

                    if s > score:
                        score, best = s, tr

                if best and score > 0.6:
                    metadata['album'] = best['album']['name']
                    metadata['year'] = best['album']['release_date'][:4] if best['album'].get('release_date') else None

                    try:
                        art = self.spotify_client.artist(best['artists'][0]['id'])
                        if art['genres']:
                            metadata['genre'] = art['genres'][0].title()
                    except:
                        pass

                    return {
                        'spotify_artist': [art['name'] for art in best['artists']],
                        'spotify_album': best['album']['name'],
                        'spotify_genres': art.get('genres', []) if 'art' in locals() else [],
                        'spotify_popularity': best['popularity'],
                        'spotify_id': best['id'],
                        'spotify_release_date': best['album'].get('release_date', '')
                    }

        except Exception as e:
            print(f"Erreur Spotify API: {e}")

        return {}

    def _calculate_match_score(self, clean_artist: str, clean_track: str, clean_album: str, track_info: Dict) -> float:
        from difflib import SequenceMatcher

        score = 0

        track_artists = [self.normalize_string(art['name']) for art in track_info['artists']]
        if clean_artist in track_artists:
            score += 50
        else:
            for art_name in track_artists:
                if art_name in clean_artist or clean_artist in art_name:
                    score += 25
                    break

        track_name = self.normalize_string(track_info['name'])
        if clean_track == track_name:
            score += 30
        else:
            similarity = SequenceMatcher(None, clean_track, track_name).ratio()
            score += similarity * 30

        if clean_album:
            album_name = self.normalize_string(track_info['album']['name'])
            if clean_album == album_name:
                score += 20
            else:
                similarity = SequenceMatcher(None, clean_album, album_name).ratio()
                score += similarity * 20

        score += track_info['popularity'] / 10

        return score

    def _extract_track_metadata(self, track_info: Dict) -> Dict[str, Any]:
        """Extrait les métadonnées d'une piste Spotify"""
        artist_id = track_info['artists'][0]['id']
        genres = self._get_artist_genres(artist_id)

        return {
            'spotify_artist': [art['name'] for art in track_info['artists']],
            'spotify_album': track_info['album']['name'],
            'spotify_genres': genres,
            'spotify_popularity': track_info['popularity'],
            'spotify_id': track_info['id'],
            'spotify_album_id': track_info['album']['id'],
            'spotify_release_date': track_info['album'].get('release_date', ''),
            'spotify_track_number': track_info.get('track_number', 0),
            'spotify_duration_ms': track_info.get('duration_ms', 0)
        }

    @staticmethod
    def normalize_string(s: str) -> str:
        """Normalise une chaîne pour la comparaison"""
        if not s:
            return ""
        s = s.lower()
        s = re.sub(r'[^\w\s]', '', s)
        s = re.sub(r'\s+', ' ', s).strip()
        return s

    def _find_track_in_artist_albums(self, artist_id: str, track_name: str) -> Dict[str, Any]:
        try:
            albums = []
            offset = 0
            limit = 50

            while True:
                results = self.spotify_client.artist_albums(
                    artist_id,
                    album_type='album,single,compilation',
                    limit=limit,
                    offset=offset
                )
                albums.extend(results['items'])

                if len(results['items']) < limit:
                    break

                offset += limit

            clean_track_name = self.normalize_string(track_name)

            for album in albums:
                try:
                    tracks = self.spotify_client.album_tracks(album['id'])['items']

                    for track in tracks:
                        if self.normalize_string(track['name']) == clean_track_name:
                            return self._extract_track_metadata(track)
                except Exception as e:
                    print(f"Erreur lors de la récupération des pistes de l'album {album['name']}: {e}")
                    continue

        except Exception as e:
            print(f"Erreur lors de la recherche dans les albums de l'artiste: {e}")

        return {}

    def _get_artist_genres(self, artist_id: str) -> List[str]:
        """Récupère les genres d'un artiste Spotify"""
        try:
            artist = self.spotify_client.artist(artist_id)
            return artist.get('genres', [])
        except:
            return []

    def on_finished(self, link: str, uploaded_file: str, category: str, title: str, is_media: bool, local_path:None):
        try:
            file_extension = os.path.splitext(uploaded_file)[1].lower()
            metadata = self.get_all_metadata(uploaded_file, category)

            if category.lower() == 'videos':
                tmdb_metadata = self.get_tmdb_metadata(uploaded_file)
                metadata['tmdb'] = tmdb_metadata
                self.tmdb_metadata[uploaded_file] = json.dumps(safe_serialize(tmdb_metadata))

            if category.lower() == 'musiques' and 'spotify' in metadata:
                self.music_metadata[uploaded_file] = json.dumps(safe_serialize(metadata['spotify']))

            new_title = title
            key = category.lower()

            if key == 'videos':
                tags = metadata.get('ffprobe', {}).get('format', {}).get('tags', {})
                if tags.get('title', '').strip():
                    new_title = tags['title'].strip()

            elif key == 'musiques':
                spotify_meta = metadata.get('spotify', {})
                basic_meta = metadata  # Métadonnées de base

                # Déterminer le meilleur titre possible
                if spotify_meta and spotify_meta.get('spotify_artist') and spotify_meta.get('spotify_album'):
                    artist = spotify_meta['spotify_artist'][0]
                    album = spotify_meta['spotify_album']

                    # Si on a un titre de piste dans les métadonnées de base
                    track_title = basic_meta.get('title', '')
                    if track_title:
                        new_title = f"{artist} - {track_title}"
                    else:
                        new_title = f"{artist} - {album}"

                # Fallback sur les métadonnées de base
                elif basic_meta.get('artist') and basic_meta.get('title'):
                    new_title = f"{basic_meta['artist']} - {basic_meta['title']}"

                # Fallback sur le nom de fichier (nettoyé)
                else:
                    # Nettoyer le titre du fichier
                    clean_title = title
                    # Supprimer les extensions
                    clean_title = re.sub(r'\.(mp3|flac|wav|m4a|aac|wma)$', '', clean_title, flags=re.IGNORECASE)
                    # Supprimer les mentions entre parenthèses
                    clean_title = re.sub(r'\([^)]*\)', '', clean_title)
                    # Supprimer les mentions entre crochets
                    clean_title = re.sub(r'\[[^\]]*\]', '', clean_title)
                    # Nettoyer les espaces multiples
                    clean_title = re.sub(r'\s+', ' ', clean_title).strip()
                    new_title = clean_title

            firebase_key_raw = sanitize_for_firebase_key(new_title).lower()
            firebase_key = sanitize_folder_name(firebase_key_raw)

            # --- Préparation du dossier cible ---
            parent_id = self.folder_ids_by_name.get(key)
            folder_id = parent_id
            if is_media or key == 'documents':
                sub_key = f"{key}/{firebase_key}"
                folder_id = self._get_or_create_folder(
                    base_key=key,
                    sub_key=sub_key,
                    name=firebase_key,
                    parent_folder_id=parent_id
                )
                if key != 'documents':
                    try:
                        self.client.move_file([link], destination_folder=folder_id)
                    except FichierResponseNotOk as e:
                        if '#604' not in str(e):
                            raise

            # --- Gestion par catégorie ---
            if key == 'images':
                # Stockage metadata sans miniature
                self.store_metadata_in_firebase(
                    main_link='',
                    metadata=metadata,
                    thumb_link='',
                    category=category,
                    title=firebase_key,
                    file_extension=file_extension
                )

                # Génération de la miniature au ratio original
                thumb_bytes, suffix = self.generate_image_thumbnail(
                    uploaded_file,
                    max_width=480,
                    max_height=480,
                    quality=95
                )
                if thumb_bytes:
                    m = MultipartEncoder({
                        'file[]': (Path(uploaded_file).stem + suffix, thumb_bytes, 'image/jpeg')
                    })
                    monitor = MultipartEncoderMonitor(m, lambda mon: None)
                    headers = {'Content-Type': monitor.content_type}
                    if self.client.authed:
                        headers.update(self.client.auth_nc)

                    resp = self.client.api_call(
                        'https://api.1fichier.com/v1/upload/get_upload_server.cgi', method='GET'
                    )
                    up_srv, upload_id = resp['url'], resp['id']
                    r = s.post(
                        f'https://{up_srv}/upload.cgi?id={upload_id}',
                        data=monitor, headers=headers, allow_redirects=False
                    )
                    loc = r.headers['Location']
                    r2 = s.get(f'https://{up_srv}{loc}')
                    link_img = re.search(
                        r'<td class="normal"><a href="(.+)"', r2.text
                    ).group(1)

                    # Déplacer miniature et mettre à jour Firebase
                    self.client.move_file([link_img], destination_folder=folder_id)
                    thumb_link = link_img
                    self._get_db_ref(category, firebase_key).update(
                        {'thumbnail_link': thumb_link}, self.token
                    )
                else:
                    # Échec miniature → on utilise le lien principal
                    thumb_link = link
                    self._get_db_ref(category, firebase_key).update(
                        {'thumbnail_link': thumb_link}, self.token
                    )

                # Émettre fin et passer au suivant
                self.files_done += 1
                self.progress.emit(int(self.files_done / self.total_files * 100))
                self.finished.emit(link, uploaded_file)
                QTimer.singleShot(0, self.upload_next_file)
                return

            self.store_metadata_in_firebase(
                main_link=link if key != 'images' else '',
                metadata=metadata,
                thumb_link='',
                category=category,
                title=firebase_key,
                file_extension=file_extension,
                local_path=uploaded_file  # Ajouter le chemin local
            )

            # Mémoriser le lien principal
            link_with_aff = link
            self._firebase_keys[uploaded_file] = (category, firebase_key, link_with_aff)
            self.existing_files.setdefault(key, set()).add(firebase_key)

            if key == 'videos':
                # Génération miniature 16:9 pour vidéo
                thumb = str(Path(uploaded_file).with_suffix('.thumb.jpg'))
                try:
                    generate_thumbnail(uploaded_file, thumb)
                    self.upload_and_move_thumb(thumb, folder_id, uploaded_file)
                except Exception as e:
                    print(f"[Thumbnail] Erreur génération vidéo : {e}")
                    self._on_thumb_finished('', uploaded_file)

            elif key == 'musiques':
                # Extraction cover audio
                cover_path = str(Path(uploaded_file).with_suffix('.thumb.jpg'))
                if extract_audio_cover(uploaded_file, cover_path):
                    self.upload_and_move_thumb(cover_path, folder_id, uploaded_file)
                    return
                # Pas de cover → on conclut
                self.files_done += 1
                self.progress.emit(int(self.files_done / self.total_files * 100))
                self.finished.emit(link_with_aff, uploaded_file)
                QTimer.singleShot(0, self.upload_next_file)

            elif key == 'documents':
                # Rendu A4 de la 1ère page
                img_bytes, suffix = self.render_first_page_a4(uploaded_file, dpi=150)
                try:
                    self.client.move_file([link], destination_folder=folder_id)
                except FichierResponseNotOk as e:
                    if '#604' not in str(e):
                        raise
                thumb_link = ''
                if img_bytes:
                    m = MultipartEncoder({
                        'file[]': (Path(uploaded_file).stem + suffix, img_bytes, 'image/jpeg')
                    })
                    monitor = MultipartEncoderMonitor(m, lambda mon: None)
                    headers = {'Content-Type': monitor.content_type}
                    if self.client.authed:
                        headers.update(self.client.auth_nc)
                    resp = self.client.api_call(
                        'https://api.1fichier.com/v1/upload/get_upload_server.cgi', method='GET'
                    )
                    up_srv, upload_id = resp['url'], resp['id']
                    r = s.post(f'https://{up_srv}/upload.cgi?id={upload_id}',
                               data=monitor, headers=headers, allow_redirects=False)
                    loc = r.headers.get('Location')
                    r2 = s.get(f'https://{up_srv}{loc}')
                    link_img = re.search(r'<td class="normal"><a href="(.+)"', r2.text).group(1)
                    thumb_link = link_img

                # Mettre à jour Firebase
                self.store_metadata_in_firebase(
                    main_link=link,
                    metadata=metadata,
                    thumb_link=thumb_link,
                    category=category,
                    title=firebase_key,
                    file_extension=file_extension
                )
                self.files_done += 1
                self.progress.emit(int(self.files_done / self.total_files * 100))
                self.finished.emit(link, uploaded_file)
                QTimer.singleShot(0, self.upload_next_file)
                return

            else:
                # Archives, exécutables, etc.
                try:
                    self.client.move_file([link], destination_folder=folder_id)
                except FichierResponseNotOk as e:
                    if '#604' not in str(e):
                        raise
                self.files_done += 1
                self.progress.emit(int(self.files_done / self.total_files * 100))
                self.finished.emit(link_with_aff, uploaded_file)
                QTimer.singleShot(0, self.upload_next_file)

        except Exception as e:
            print(f"[UploadManager] Erreur sur {uploaded_file}: {e}")
            self.files_done += 1
            self.progress.emit(int(self.files_done / self.total_files * 100))
            self.finished.emit(link, uploaded_file)
            QTimer.singleShot(0, self.upload_next_file)

    @staticmethod
    def generate_image_thumbnail(input_path: str,
                                 max_width: int = 480,
                                 max_height: int = 480,
                                 quality: int = 95) -> tuple[bytes, str] | tuple[None, None]:
        try:
            with Image.open(input_path) as img:
                # Conversion en RGB si nécessaire
                if img.mode == 'RGBA':
                    bg = Image.new('RGB', img.size, (255, 255, 255))
                    bg.paste(img, mask=img.split()[3])
                    img = bg
                elif img.mode != 'RGB':
                    img = img.convert('RGB')

                # Redimensionnement (ratio conservé)
                img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)

                # Sauvegarde en JPEG sans sous‑échantillonnage (subsampling=0)
                buf = BytesIO()
                img.save(buf, 'JPEG',
                         quality=quality,
                         optimize=True,
                         subsampling=0)
                return buf.getvalue(), '.jpg'
        except Exception as e:
            print(f"[Thumbnail Images] Erreur génération miniature : {e}")
            return None, None

    def _get_or_create_folder(self, base_key: str, sub_key: str, name: str, parent_folder_id: str) -> str:
        if sub_key in self.folder_ids_by_name:
            return self.folder_ids_by_name[sub_key]

        try:
            resp = self.client.create_folder(name, parent_folder_id=parent_folder_id)
            folder_id = resp.get('folder_id') or resp.get('id')
            self.folder_ids_by_name[sub_key] = folder_id
            return folder_id
        except Exception as e:
            if "Folder already exist" in str(e):
                # Le dossier existe déjà, on doit le retrouver
                try:
                    # Lister les dossiers dans le parent pour trouver celui qui existe
                    parent_folders = self.client.get_folders(parent_folder_id)
                    folders = parent_folders.get('sub_folders', []) if isinstance(parent_folders, dict) else []

                    for folder in folders:
                        folder_name = folder.get('name', '').lower() if isinstance(folder,
                                                                                   dict) else folder.name.lower()
                        if folder_name == name.lower():
                            folder_id = folder.get('id') if isinstance(folder, dict) else folder.id
                            self.folder_ids_by_name[sub_key] = folder_id
                            return folder_id

                    # Si on arrive ici, on n'a pas trouvé le dossier, on utilise le parent comme fallback
                    print(f"Dossier '{name}' existe mais non trouvé, utilisation du parent")
                    return parent_folder_id
                except Exception as e2:
                    print(f"Erreur lors de la recherche du dossier existant: {e2}")
                    return parent_folder_id
            else:
                raise

    def _complete_document_upload(self, link: str, uploaded_file: str):
        self.files_done += 1
        self.progress.emit(int(self.files_done / self.total_files * 100))
        self.finished.emit(link, uploaded_file)
        QTimer.singleShot(0, self.upload_next_file)

    @staticmethod
    def extract_first_image_bytes(doc_path: str) -> tuple[bytes, str] | tuple[None, None]:
        path = Path(doc_path)
        ext = path.suffix.lower()

        if ext == '.pdf':
            doc = fitz.open(str(path))
            if len(doc) > 0:
                pix = doc.load_page(0).get_pixmap()
                data = pix.tobytes('png')
                return data, '.png'
            doc.close()

        elif ext == '.docx':
            doc = Document(str(path))
            for rel in doc.part._rels.values():
                if 'image' in rel.reltype:
                    data = rel.target_part.blob
                    suffix = Path(rel.target_part.partname).suffix
                    return data, suffix

        elif ext == '.pptx':
            pres = Presentation(str(path))
            if pres.slides:
                for shape in pres.slides[0].shapes:
                    if hasattr(shape, 'image'):
                        data = shape.image.blob
                        suffix = f'.{shape.image.ext}'
                        return data, suffix

        return None, None

    @staticmethod
    def resize_image_thumbnail(input_path: str, output_path: str, target_size=(640, 360)) -> bool:
        try:

            with Image.open(input_path) as img:
                if img.mode == 'RGBA':
                    bg = Image.new('RGB', img.size, (0, 0, 0))
                    bg.paste(img, mask=img.split()[3])
                    img = bg
                elif img.mode != 'RGB':
                    img = img.convert('RGB')

                img.thumbnail(target_size, Image.Resampling.LANCZOS)

                thumb = Image.new('RGB', target_size, (0, 0, 0))
                x = (target_size[0] - img.width) // 2
                y = (target_size[1] - img.height) // 2
                thumb.paste(img, (x, y))

            thumb.save(output_path, 'JPEG',
                       quality=100,
                       subsampling=0,
                       optimize=False)

            return True

        except Exception as e:
            print(f"[Thumbnail] Erreur redimensionnement image : {e}")
            return False

    def upload_and_move_thumb(self, thumb_path: str, dest_folder_id: str, parent_file: str):
        t = UploadThread(self.client, thumb_path)
        t.progress.connect(lambda _: None)
        t.finished.connect(lambda link: self.client.move_file([link], destination_folder=dest_folder_id))
        t.finished.connect(lambda link, pf=parent_file: self.thumb_finished.emit(link, pf))
        t.finished.connect(lambda: os.remove(thumb_path))
        self.thumb_threads.append(t)
        t.start()

    def _on_thumb_finished(self, thumb_url: str, parent_file: str):
        self._ensure_token()
        category, title, main_link = self._firebase_keys[parent_file]
        thumb_url = thumb_url
        try:
            self._get_db_ref(category, title).update({'thumbnail_link': thumb_url}, self.token)
        except Exception as e:
            print(f"[Firebase] Erreur update thumbnail: {e}")
        self.thumb_threads = [t for t in self.thumb_threads if t.isRunning()]
        self.files_done += 1
        self.progress.emit(int(self.files_done / self.total_files * 100))
        self.finished.emit(main_link, parent_file)
        QTimer.singleShot(0, self.upload_next_file)

    @staticmethod
    def get_image_basic_info(file_path: str) -> dict:
        """
        Fallback pour images corrompues - lit les infos de base du fichier
        """
        path = Path(file_path)
        metadata = {
            'format': path.suffix.lstrip('.').upper(),
            'file_size': path.stat().st_size if path.exists() else 0,
            'size': {}
        }

        try:
            with open(file_path, 'rb') as f:
                header = f.read(16)

            # Détecter format par magic bytes
            if header.startswith(b'\xff\xd8\xff'):
                metadata['format'] = 'JPEG'
            elif header.startswith(b'\x89PNG\r\n\x1a\n'):
                metadata['format'] = 'PNG'
            elif header.startswith(b'GIF87a') or header.startswith(b'GIF89a'):
                metadata['format'] = 'GIF'
            elif header.startswith(b'BM'):
                metadata['format'] = 'BMP'

        except Exception:
            pass

        return metadata

    def get_all_metadata(self, file_path: str, category: str) -> dict:
        """
        Extrait les métadonnées selon la catégorie avec fallback robuste pour les images
        """
        metadata = {}
        try:
            key = category.lower()

            if key == 'videos':
                proc = subprocess.run(
                    [str(FFPROBE_PATH),
                     '-v', 'quiet',
                     '-print_format', 'json',
                     '-show_format',
                     '-show_streams',
                     file_path],
                    capture_output=True,
                    text=True,
                    check=True
                )
                metadata['ffprobe'] = json.loads(proc.stdout)

            elif key == 'musiques':
                # Métadonnées de base avec la nouvelle fonction
                metadata.update(extract_music_metadata(file_path))

                # Métadonnées enrichies avec Spotify
                if self.spotify_client and metadata.get('artist'):
                    spotify_meta = self.get_spotify_metadata(
                        artist=metadata['artist'],
                        track=metadata.get('title'),
                        album=metadata.get('album')
                    )
                    if spotify_meta:
                        metadata['spotify'] = spotify_meta

            elif key == 'images':
                try:
                    # Essayer PIL d'abord
                    with Image.open(file_path) as img:
                        metadata['format'] = img.format
                        metadata['size'] = {'width': img.width, 'height': img.height}
                        exif_data = img._getexif() or {}
                        exif = {}
                        for tag_id, value in exif_data.items():
                            tag = ExifTags.TAGS.get(tag_id, tag_id)
                            exif[tag] = value
                        if exif:
                            metadata['exif'] = exif
                except Exception:
                    # Fallback simple - juste les infos de fichier
                    print(f"[Metadata] PIL échoue sur {file_path}, utilisation du fallback")
                    metadata.update(UploadManager.get_image_basic_info(file_path))

        except Exception as e:
            print(f"[Metadata] Erreur lecture métadonnées : {e}")

        return metadata